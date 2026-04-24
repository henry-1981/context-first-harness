#!/usr/bin/env python3
import json
import sys
from html import escape

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({"error": "python-pptx is not installed"}))
    sys.exit(2)


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "usage: pptx_parser.py <file> [slide_index]"}))
        sys.exit(1)

    path = sys.argv[1]
    slide_index = int(sys.argv[2]) if len(sys.argv) > 2 else 0
    prs = Presentation(path)
    if slide_index >= len(prs.slides):
      print(json.dumps({"error": "slide index out of range"}))
      sys.exit(1)

    slide = prs.slides[slide_index]
    blocks = []
    html_parts = ['<div data-source-kind="pptx">']
    tables = []
    slide_size = {
        "width": int(prs.slide_width),
        "height": int(prs.slide_height),
    }
    order = 0

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                block = {
                    "type": "text",
                    "role": "title" if "Title" in getattr(shape, "name", "") else "text",
                    "text": text,
                    "order": order,
                    "bbox": {
                        "x": int(shape.left),
                        "y": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    },
                }
                blocks.append(block)
                html_parts.append(f'<div data-block-type="text">{escape(text)}</div>')
                order += 1
        elif getattr(shape, "shape_type", None) == 13:
            block = {
                "type": "image",
                "role": "image",
                "text": None,
                "order": order,
                "bbox": {
                    "x": int(shape.left),
                    "y": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                },
            }
            blocks.append(block)
            html_parts.append('<div data-block-type="image"></div>')
            order += 1
        elif getattr(shape, "has_table", False):
            rows = []
            table_html = ['<table data-block-type="table">']
            for r_idx, row in enumerate(shape.table.rows):
                row_cells = []
                table_html.append('<tr>')
                for c_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    row_cells.append({
                        "row": r_idx,
                        "col": c_idx,
                        "text": cell_text,
                        "bbox": {
                            "x": int(shape.left + c_idx * (shape.width / max(len(row.cells), 1))),
                            "y": int(shape.top + r_idx * (shape.height / max(len(shape.table.rows), 1))),
                            "width": int(shape.width / max(len(row.cells), 1)),
                            "height": int(shape.height / max(len(shape.table.rows), 1)),
                        },
                        "rowspan": 1,
                        "colspan": 1,
                    })
                    table_html.append(f'<td>{escape(cell_text)}</td>')
                table_html.append('</tr>')
                rows.append({"cells": row_cells})
            table_html.append('</table>')
            table = {
                "type": "table",
                "role": "table",
                "order": order,
                "rows": rows,
                "bbox": {
                    "x": int(shape.left),
                    "y": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                },
            }
            tables.append(table)
            blocks.append(table)
            html_parts.append(''.join(table_html))
            order += 1

    html_parts.append('</div>')

    print(json.dumps({
        "html_string": ''.join(html_parts),
        "metadata": {
            "source_kind": "pptx",
            "slide_count": len(prs.slides),
            "slide_size": slide_size,
            "warnings": [],
        },
        "extraction_report": {
            "blocks": blocks,
            "tables": tables,
            "confidence": 0.7,
            "fallback_used": False,
        },
    }))


if __name__ == "__main__":
    main()
